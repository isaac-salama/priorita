#!/usr/bin/env python3
from pptx import Presentation
import html

def extract_ppt_content(ppt_path):
    """Extract content from PowerPoint file"""
    prs = Presentation(ppt_path)
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        slide_content = {
            'title': '',
            'text': [],
            'notes': ''
        }
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                # Check if it's a title (usually first shape or has title formatting)
                if not slide_content['title'] and len(text) < 100:
                    slide_content['title'] = text
                else:
                    slide_content['text'].append(text)
        
        # Extract notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                slide_content['notes'] = notes_slide.notes_text_frame.text
        
        slides_data.append(slide_content)
    
    return slides_data

def create_reveal_html(slides_data, output_path):
    """Create Reveal.js HTML from slides data"""
    
    html_content = '''<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Finca La Priorita - Dossier 2022</title>
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.3.1/dist/reveal.css">
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.3.1/dist/theme/white.css">
    <style>
        .reveal .slides section {
            text-align: left;
        }
        .reveal h1, .reveal h2, .reveal h3 {
            color: #2c3e50;
        }
        .reveal .slides section .fragment {
            opacity: 0;
        }
        .reveal .slides section .fragment.visible {
            opacity: 1;
        }
        body {
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
        }
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">
'''
    
    for slide in slides_data:
        html_content += '            <section>\n'
        
        if slide['title']:
            html_content += f'                <h2>{html.escape(slide["title"])}</h2>\n'
        
        if slide['text']:
            for text in slide['text']:
                if text and text != slide['title']:
                    # Split long text into paragraphs
                    paragraphs = text.split('\n')
                    for para in paragraphs:
                        if para.strip():
                            html_content += f'                <p>{html.escape(para.strip())}</p>\n'
        
        html_content += '            </section>\n'
    
    html_content += '''        </div>
    </div>
    
    <script src="https://cdn.jsdelivr.net/npm/reveal.js@4.3.1/dist/reveal.js"></script>
    <script>
        Reveal.initialize({
            hash: true,
            controls: true,
            progress: true,
            center: true,
            touch: true,
            transition: 'slide'
        });
    </script>
</body>
</html>'''
    
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(html_content)
    
    print(f"Created {output_path} with {len(slides_data)} slides")

if __name__ == "__main__":
    ppt_file = "DOSSIER FINCA LA PRIORITA 2022.ppt"
    output_file = "dossier.html"
    
    print(f"Extracting content from {ppt_file}...")
    slides_data = extract_ppt_content(ppt_file)
    print(f"Found {len(slides_data)} slides")
    
    print(f"Creating web slideshow: {output_file}...")
    create_reveal_html(slides_data, output_file)
    print("Done!")

